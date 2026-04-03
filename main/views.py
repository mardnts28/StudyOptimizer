from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.models import User
import os
try:
    from decouple import config
except ImportError:
    config = os.environ.get

from django.contrib import messages
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.views.decorators.csrf import csrf_protect
from django.views.decorators.http import require_POST
from django.http import JsonResponse, HttpResponse
from django.db.models import Sum, Count
from datetime import date, timedelta, datetime
import json
import re
import bleach
import io
import PyPDF2
import docx
from pptx import Presentation
from collections import Counter
import random
from django.core.mail import send_mail
from django.conf import settings
from .models import Task, SharedMaterial, Comment, SummarizedDocument, ScheduleItem, UserProfile, AuditLog
import pyotp
import qrcode
import qrcode.image.svg
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests
import google.generativeai as genai
from django_ratelimit.decorators import ratelimit
from .utils import send_security_alert
from .models import KnownIP


@csrf_protect
def index(request):
    return render(request, "main/index.html")


@csrf_protect
def register(request):
    if request.user.is_authenticated:
        return redirect("dashboard")
    if request.method == "POST":
        # Honeypot check (SPAM protection)
        if request.POST.get('_hp_field'):
            # Return a cryptic message or 200 to confuse bots, but don't proceed
            return HttpResponse("Registration successful! (Not really - bot detected)", status=200)

        username = request.POST.get('username', '').strip()
        email    = request.POST.get('email', '').strip()
        password = request.POST.get('password', '')
        password2 = request.POST.get('password2', '')

        from django.core.validators import validate_email
        from django.core.exceptions import ValidationError
        from django.contrib.auth.password_validation import validate_password

        if not re.match(r'^[a-zA-Z0-9_\.\-]{3,150}$', username):
            messages.error(request, "Invalid username format.")
            return redirect('register')

        try:
            validate_email(email)
        except ValidationError:
            messages.error(request, "Invalid email format.")
            return redirect('register')

        if password != password2:
            messages.error(request, "Passwords do not match.")
            return redirect('register')

        try:
            temp_user = User(username=username, email=email)
            validate_password(password, temp_user)
        except ValidationError as e:
            for message in e.messages:
                messages.error(request, message)
            return redirect('register')

        if User.objects.filter(username=username).exists():
            messages.error(request, "Username already taken.")
            return redirect('register')
        if User.objects.filter(email=email).exists():
            messages.error(request, "Email is already registered.")
            return redirect('register')

        User.objects.create_user(username=username, email=email, password=password)
        messages.success(request, "Account created successfully. Please log in.")
        return redirect('login')

    return render(request, 'main/register.html')


def start_mfa(request, user):
    profile, _ = UserProfile.objects.get_or_create(user=user)
    request.session['mfa_user_id'] = user.id
    if profile.totp_enabled:
        request.session['mfa_method'] = 'totp'
        return redirect('mfa_verify')
    else:
        return redirect('setup_totp')

def login_view(request):
    if request.user.is_authenticated:
        return redirect("dashboard")
    if request.method == "POST":
        email    = request.POST.get("email")
        password = request.POST.get("password")
        try:
            user_obj = User.objects.get(email=email)
            user = authenticate(request, username=user_obj.username, password=password)
        except User.DoesNotExist:
            user = None

        if user is not None:
            return start_mfa(request, user)
        else:
            messages.error(request, "Invalid email or password.")

    return render(request, "main/login.html")

@csrf_protect
@require_POST
def google_login(request):
    from django.http import Http404
    raise Http404("Google Sign-In is temporarily disabled.")

def setup_totp(request):
    if request.user.is_authenticated:
        return redirect("dashboard")

    mfa_user_id = request.session.get('mfa_user_id')
    if not mfa_user_id:
        return redirect('login')

    user = get_object_or_404(User, id=mfa_user_id)
    profile, _ = UserProfile.objects.get_or_create(user=user)

    if profile.totp_enabled:
        return redirect('mfa_verify')

    secret = profile.generate_totp_secret()
    
    totp = pyotp.TOTP(secret)
    provisioning_uri = totp.provisioning_uri(name=user.email, issuer_name="StudyOptimizer")
    
    factory = qrcode.image.svg.SvgPathImage
    img = qrcode.make(provisioning_uri, image_factory=factory)
    stream = io.BytesIO()
    img.save(stream)
    svg_data = stream.getvalue().decode()

    if request.method == "POST":
        entered_code = request.POST.get("otp", "").strip()
        if pyotp.TOTP(secret).verify(entered_code):
            profile.totp_enabled = True
            profile.save()
            
            from django.contrib.auth.backends import ModelBackend
            user.backend = 'django.contrib.auth.backends.ModelBackend'
            login(request, user)
            del request.session['mfa_user_id']
            return redirect("dashboard")
        else:
            messages.error(request, "Invalid code. Please try scanning again.")

    return render(request, "main/setup_totp.html", {'qr_code': svg_data, 'secret': secret})

def mfa_verify(request):
    if request.user.is_authenticated:
        return redirect("dashboard")

    mfa_user_id = request.session.get('mfa_user_id')
    mfa_method = request.session.get('mfa_method', 'totp')

    if not mfa_user_id:
        messages.error(request, "Verification session expired. Please log in again.")
        return redirect('login')

    user = get_object_or_404(User, id=mfa_user_id)
    profile, _ = UserProfile.objects.get_or_create(user=user)

    if request.method == "POST":
        action = request.POST.get('action')
        
        if action == 'send_email':
            otp = f"{random.randint(100000, 999999)}"
            request.session['mfa_otp'] = otp
            request.session['mfa_method'] = 'email'
            
            try:
                send_mail(
                    'Your Study Optimizer Verification Code',
                    f'Your verification code is: {otp}\n\nPlease enter this code to securely log in.',
                    settings.DEFAULT_FROM_EMAIL,
                    [user.email],
                    fail_silently=False,
                )
                messages.info(request, "A verification code has been sent to your email.")
            except Exception as e:
                messages.error(request, f"Failed to send email: {e}")
                
            return redirect('mfa_verify')

        entered_code = request.POST.get("otp", "").strip()
        
        if mfa_method == 'email':
            expected_otp = request.session.get('mfa_otp')
            if expected_otp and entered_code == str(expected_otp):
                from django.contrib.auth.backends import ModelBackend
                user.backend = 'django.contrib.auth.backends.ModelBackend'
                login(request, user)
                
                del request.session['mfa_user_id']
                if 'mfa_otp' in request.session: del request.session['mfa_otp']
                if 'mfa_method' in request.session: del request.session['mfa_method']
                return redirect("dashboard")
            else:
                messages.error(request, "Invalid verification code.")
        else:
            totp = pyotp.TOTP(profile.totp_secret)
            if totp.verify(entered_code):
                from django.contrib.auth.backends import ModelBackend
                user.backend = 'django.contrib.auth.backends.ModelBackend'
                
                # --- STRIDE Enhanced: New IP Detection ---
                current_ip = request.META.get('REMOTE_ADDR', '127.0.0.1')
                if not KnownIP.objects.filter(user=user, ip_address=current_ip).exists():
                    send_security_alert(
                        user, 
                        "New Login Device Detected", 
                        f"Your account was just logged into from a new device or location.\n\nIP Address: {current_ip}\nDate: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                    )
                    KnownIP.objects.create(user=user, ip_address=current_ip)
                else:
                    # Update last used
                    KnownIP.objects.filter(user=user, ip_address=current_ip).update(last_used=datetime.now())

                login(request, user)
                
                del request.session['mfa_user_id']
                if 'mfa_method' in request.session: del request.session['mfa_method']
                return redirect("dashboard")
            else:
                messages.error(request, "Invalid authenticator code.")

    return render(request, "main/mfa_verify.html", {'mfa_method': mfa_method, 'email': user.email})

def logout_view(request):
    logout(request)
    return redirect("login")



@login_required
@csrf_protect
def dashboard(request):
    user_tasks     = Task.objects.filter(user=request.user)
    active_tasks   = user_tasks.filter(completed=False).order_by('due_date')
    total_tasks    = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    completion_rate = round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0

    upcoming_tasks_list = []
    for t in active_tasks[:4]:
        days_left = (t.due_date - date.today()).days
        upcoming_tasks_list.append({
            'title':    t.title,
            'date':     t.due_date.strftime('%b %d'),
            'priority': t.priority,
            'category':   'General',
            'daysLeft': max(days_left, 0),
        })

    today          = date.today()
    start_of_week  = today - timedelta(days=today.weekday())
    daily_hours    = []
    for i in range(7):
        day = start_of_week + timedelta(days=i)
        tasks_on_day    = user_tasks.filter(completed=True, created_at__date=day).count()
        summaries_on_day = SummarizedDocument.objects.filter(user=request.user, created_at__date=day).count()
        daily_hours.append((tasks_on_day * 2) + summaries_on_day)

    recent_summaries = SummarizedDocument.objects.filter(user=request.user).order_by('-created_at')[:5]
    summaries_list   = [{
        'id':      s.id,
        'title':   s.file_name,
        'date':    s.created_at.strftime('%b %d'),
        'emoji':   s.emoji,
        'summary': s.summary_text[:100] + '...' if len(s.summary_text) > 100 else s.summary_text,
    } for s in recent_summaries]

    context = {
        'total_tasks':        total_tasks,
        'completed_count':    completed_tasks,
        'completion_rate':    completion_rate,
        'upcoming_tasks_json': json.dumps(upcoming_tasks_list),
        'recent_summaries_json': json.dumps(summaries_list),
        'docs_count':         SharedMaterial.objects.filter(author=request.user).count()
                              + SummarizedDocument.objects.filter(user=request.user).count(),
        'study_hours':        (completed_tasks * 2) + SummarizedDocument.objects.filter(user=request.user).count(),
        'weekly_hours_list':  json.dumps(daily_hours),
        'schedule_items_json': json.dumps([{
            'id':       item.id,
            'day':      item.day,
            'time':     item.time,
            'activity': item.activity,
            'color':    item.color,
        } for item in ScheduleItem.objects.filter(user=request.user)]),
    }
    return render(request, "main/dashboard.html", context)



@login_required
def tasks_view(request):
    tasks      = Task.objects.filter(user=request.user)
    tasks_data = []
    for t in tasks:
        tasks_data.append({
            'id':        t.id,
            'title':     t.title,
            'subject':   t.subject,
            'category':    t.category,
            'priority':  t.priority,
            'dueDate':   t.due_date.strftime('%Y-%m-%d'),
            'completed': t.completed,
        })
    return render(request, 'main/tasks.html', {'tasks_data': json.dumps(tasks_data)})


@login_required
@require_POST
@ratelimit(key='ip', rate='10/m', block=True)
def add_task(request):
    try:
        data     = json.loads(request.body)
        
        title = data.get('title', '').strip()
        if not title or len(title) > 255:
            return JsonResponse({'status': 'error', 'message': 'Invalid title length.'}, status=400)
            
        priority = data.get('priority')
        if priority not in ['Low', 'Medium', 'High']:
            return JsonResponse({'status': 'error', 'message': 'Invalid priority.'}, status=400)
            
        due_date_str = data.get('dueDate', '')
        if not re.match(r'^\d{4}-\d{2}-\d{2}$', due_date_str):
            return JsonResponse({'status': 'error', 'message': 'Invalid due date structure.'}, status=400)

        due_date = datetime.strptime(due_date_str, '%Y-%m-%d').date()

        task     = Task.objects.create(
            user      = request.user,
            title     = title,
            subject   = str(data.get('subject', 'General'))[:100],
            category  = str(data.get('category', 'General'))[:20],
            priority  = priority,
            due_date  = due_date,
            completed = False,
        )
        return JsonResponse({'status': 'success', 'task': {
            'id':        task.id,
            'title':     task.title,
            'subject':   task.subject,
            'category':    task.category,
            'priority':  task.priority,
            'dueDate':   task.due_date.strftime('%Y-%m-%d'),
            'completed': task.completed,
        }})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@login_required
@require_POST
def edit_task(request, task_id):
    try:
        data     = json.loads(request.body)
        task     = Task.objects.get(id=task_id, user=request.user)
        
        title = data.get('title', '').strip()
        if not title or len(title) > 255:
            return JsonResponse({'status': 'error', 'message': 'Invalid title length.'}, status=400)
            
        priority = data.get('priority')
        if priority not in ['Low', 'Medium', 'High']:
            return JsonResponse({'status': 'error', 'message': 'Invalid priority.'}, status=400)
            
        due_date_str = data.get('dueDate', '')
        if not re.match(r'^\d{4}-\d{2}-\d{2}$', due_date_str):
            return JsonResponse({'status': 'error', 'message': 'Invalid date format.'}, status=400)
            
        due_date = datetime.strptime(due_date_str, '%Y-%m-%d').date()

        task.title    = title
        task.subject  = str(data.get('subject', task.subject))[:100]
        task.category   = str(data.get('category', task.category))[:20]
        task.priority = priority
        task.due_date = due_date
        task.save()
        return JsonResponse({'status': 'success', 'task': {
            'id':        task.id,
            'title':     task.title,
            'subject':   task.subject,
            'category':    task.category,
            'priority':  task.priority,
            'dueDate':   task.due_date.strftime('%Y-%m-%d'),
            'completed': task.completed,
        }})
    except Task.DoesNotExist:
        return JsonResponse({'status': 'error', 'message': 'Task not found'}, status=404)
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@login_required
@require_POST
def delete_task(request, task_id):
    try:
        Task.objects.get(id=task_id, user=request.user).delete()
        return JsonResponse({'status': 'success'})
    except Task.DoesNotExist:
        return JsonResponse({'status': 'error', 'message': 'Task not found'}, status=404)


@login_required
@require_POST
def toggle_task(request, task_id):
    try:
        task           = Task.objects.get(id=task_id, user=request.user)
        task.completed = not task.completed
        task.save()
        return JsonResponse({'status': 'success', 'completed': task.completed})
    except Task.DoesNotExist:
        return JsonResponse({'status': 'error', 'message': 'Task not found'}, status=404)



@login_required
@csrf_protect
def upload(request):
    recent_summaries = SummarizedDocument.objects.filter(user=request.user).order_by('-created_at')[:10]
    summaries_list   = [{
        'id':      s.id,
        'title':   s.file_name,
        'category':  'General',
        'date':    s.created_at.strftime('%b %d'),
        'emoji':   s.emoji,
        'summary': s.summary_text,
    } for s in recent_summaries]
    return render(request, "main/upload.html", {'recent_summaries': summaries_list})


@login_required
@csrf_protect
@ratelimit(key='user', rate='20/m', method='POST', block=True)
def summarize_doc(request):
    try:
        if 'file' not in request.FILES:
            return JsonResponse({'status': 'error', 'message': 'No file uploaded'}, status=400)

        uploaded_file = request.FILES['file']
        file_name     = str(uploaded_file.name).strip()

        if uploaded_file.size > 15 * 1024 * 1024:
            return JsonResponse({'status': 'error', 'message': 'File size exceeds 15MB limit.'}, status=400)
        if not file_name or len(file_name) > 255:
            return JsonResponse({'status': 'error', 'message': 'File name invalid or too long.'}, status=400)
            
        allowed_exts = ('.pdf', '.docx', '.pptx', '.txt')
        if not file_name.lower().endswith(allowed_exts):
            return JsonResponse({'status': 'error', 'message': 'Unsupported file type.'}, status=400)

        import hashlib
        uploaded_file.seek(0)
        file_hash = hashlib.sha256(uploaded_file.read()).hexdigest()
        uploaded_file.seek(0)
            
        content       = ""

        if file_name.lower().endswith('.pdf'):
            import PyPDF2
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages[:15]:
                text = page.extract_text()
                if text:
                    content += text + " "

        elif file_name.lower().endswith('.docx'):
            import docx
            doc     = docx.Document(uploaded_file)
            content = " ".join([p.text for p in doc.paragraphs[:50]])

        elif file_name.lower().endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += str(shape.text) + " "

        elif file_name.lower().endswith('.txt'):
            content = uploaded_file.read().decode('utf-8')

        else:
            return JsonResponse({'status': 'error', 'message': 'Unsupported file type.'}, status=400)

        if not content.strip():
            return JsonResponse({'status': 'error', 'message': 'Could not extract text from document.'}, status=400)
            
        content_clean = re.sub(r'\s+', ' ', content).strip()

        # --- STRIDE Enhanced AI Summarizer (Gemini Integration) ---
        final_summary = "AI Summarization is currently unavailable. Please check your document content."
        title_line    = f"📑 Document Summary – {file_name}"
        success       = False
        
        gemini_key = config('GOOGLE_API_KEY', default=None)
        if gemini_key:
            try:
                print(f"DEBUG: Initializing Gemini (REST) with key: {gemini_key[:5]}...")
                genai.configure(api_key=gemini_key, transport='rest')
                
                prompt = (
                    "Summarize this document in a detailed reviewer format. Follow this exact structure:\n"
                    "1. **Title**: The main subject or theme.\n"
                    "2. **Purpose / Objective**: Why it exists and its aims.\n"
                    "3. **Highlights**: 5–7 key points (detailed for study).\n"
                    "4. **Detailed Breakdown**: Section-by-section notes, capturing names, technical terms, and cause-effect relationships.\n"
                    "5. **Key Facts**: Bullet-point critical details (dates, figures, technical stacks).\n"
                    "6. **Cause–Effect / Progression**: Show how events/decisions led to others.\n"
                    "7. **Takeaways / Lessons**: 3–5 insights for quick recall.\n"
                    "8. **Closing Insight**: One sentence on impact or legacy.\n\n"
                    f"Content: {content_clean[:30000]}"
                )

                # Multi-model Fallback Strategy
                for model_name in ['gemini-1.5-flash', 'gemini-1.5-flash-latest', 'gemini-1.5-pro', 'gemini-1.5-pro-latest', 'gemini-1.0-pro']:
                    try:
                        print(f"DEBUG: Attempting AI with model: {model_name}")
                        model = genai.GenerativeModel(
                            model_name,
                            safety_settings=[
                                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
                            ]
                        )
                        response = model.generate_content(prompt)
                        if hasattr(response, 'text') and response.text.strip():
                            final_summary = response.text
                            title_line = f"📘 AI Generated – {file_name}"
                            success = True
                            break
                    except Exception as model_err:
                        print(f"DEBUG: Model {model_name} failed: {str(model_err)}")
                        continue
            except Exception as e:
                print(f"DEBUG: Summarize_Doc AI Error: {str(e)}")

        # Fallback to text extract if AI was not used or failed
        if not success:
            sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', content_clean) if len(s.strip()) > 10]
            if sentences:
                final_summary = "**📘 Text Extract Summary:**\n\n" + "\n".join(f"• {s}" for s in sentences[:8]) + "\n\n🌟 Success! You're making great progress. Review these key concepts to solidify your understanding and stay on track for your study goals!"
                title_line    = f"📑 Document Summary (Extraction) – {file_name}"
            else:
                final_summary = "Could not extract sufficient text for a summary. Please verify the document content."

        # Convert Markdown Bold to HTML Bold for web display
        final_summary = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', final_summary)

        doc = SummarizedDocument.objects.create(
            user         = request.user,
            file_name    = bleach.clean(file_name),
            category     = 'General',
            summary_text = final_summary,
            content_hash = hashlib.sha256(final_summary.encode('utf-8')).hexdigest(),
            emoji        = '📄',
        )
        # Audit Security
        log = AuditLog(user=request.user, action="Summarized Document")
        log.details = f"File: {file_name}, Hash: {file_hash}"
        log.save()
        return JsonResponse({'status': 'success', 'summary': final_summary, 'title': title_line, 'doc_id': doc.id})

    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)


@login_required
@require_POST
def summarize_batch(request):
    try:
        data    = json.loads(request.body)
        doc_ids = data.get('doc_ids', [])

        if not doc_ids:
            return JsonResponse({'status': 'error', 'message': 'No documents provided.'}, status=400)

        docs = SummarizedDocument.objects.filter(id__in=doc_ids, user=request.user)
        if not docs.exists():
            return JsonResponse({'status': 'error', 'message': 'Documents not found.'}, status=404)

        all_text_full = "\n---\n".join(d.summary_text for d in docs)
        
        # --- STRIDE Enhanced: Batch AI Synthesis ---
        gemini_key = config('GOOGLE_API_KEY', default=None)
        if gemini_key:
            try:
                print(f"DEBUG: Initializing Batch AI (REST) with key: {gemini_key[:5]}...")
                genai.configure(api_key=gemini_key, transport='rest')

                prompt = (
                    "Synthesize these documents into a single reviewer format. Follow this exact structure:\n"
                    "1. **Title**: Combined subject or theme.\n"
                    "2. **Purpose / Objective**: Collective aims.\n"
                    "3. **Highlights**: 5–7 integrated points.\n"
                    "4. **Detailed Breakdown**: Consolidated flow across all docs.\n"
                    "5. **Key Facts**: Critical combined details.\n"
                    "6. **Cause–Effect / Progression**: Integrated relationships.\n"
                    "7. **Takeaways / Lessons**: Strategic insights.\n"
                    "8. **Closing Insight**: Collective impact.\n\n"
                    "Content: " + all_text_full[:30000]
                )
                
                print("DEBUG: Listing available Gemini models...")
                try:
                    for m in genai.list_models():
                        if 'generateContent' in m.supported_generation_methods:
                            print(f"DEBUG: Model found: {m.name}")
                except Exception as list_err:
                    print(f"DEBUG: Could not list models: {str(list_err)}")

                # Multi-model Fallback Strategy for Batch
                success = False
                for model_name in ['gemini-1.5-flash', 'gemini-1.5-flash-latest', 'gemini-1.5-pro', 'gemini-1.5-pro-latest', 'gemini-1.0-pro']:
                    try:
                        print(f"DEBUG: Attempting Batch AI with model: {model_name}")
                        model = genai.GenerativeModel(
                            model_name,
                            safety_settings=[
                                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
                            ]
                        )
                        response = model.generate_content(prompt)
                        if hasattr(response, 'text') and response.text.strip():
                            batch_output = response.text
                            success = True
                            break
                    except Exception as model_err:
                        print(f"DEBUG: Batch Model {model_name} failed: {str(model_err)}")
                        continue
                
                if not success:
                    raise Exception("All Batch AI models failed or were blocked. (v1beta/v1 mismatch or safety)")
                    
            except Exception as e:
                print(f"DEBUG: Batch AI Synthesis Error: {str(e)}")

        if not success:
            batch_output = "Master AI Synthesis is momentarily unavailable, but don't let that slow you down! You've already processed these files—review the individual summaries below, and keep crushing your study goals! 🚀"

        # Construct full batch output including individual summaries for accuracy in PDF
        individual_summaries = ""
        for i, d in enumerate(docs, 1):
            individual_summaries += f"**📄 Document {i}: {d.file_name}**\n{d.summary_text}\n\n"
        
        full_batch_summary = f"{individual_summaries}**🔄 --- Master AI Synthesis ---**\n\n{batch_output}"

        # Convert Markdown Bold to HTML Bold for web display
        batch_output = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', batch_output)
        # Also clean individual previews if needed
        full_batch_summary = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', full_batch_summary)

        # Save the batch summary as a new document so it can be downloaded/viewed
        batch_doc = SummarizedDocument.objects.create(
            user         = request.user,
            file_name    = f"Batch Summary ({docs.count()} Files)",
            summary_text = full_batch_summary,
            emoji        = '📊',
            category     = 'Batch'
        )

        return JsonResponse({
            'status': 'success', 
            'combined_summary': batch_output, 
            'full_summary': full_batch_summary,
            'batch_doc_id': batch_doc.id,
            'batch_title': batch_doc.file_name
        })

    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)



@login_required
@csrf_protect
def collaborate(request):
    materials             = SharedMaterial.objects.all().order_by('-created_at')
    total_community_likes = 0
    materials_list        = []

    for m in materials:
        m_likes = m.likes.count()
        total_community_likes += m_likes
        
        is_anon = m.is_anonymous

        # If it's anonymous, hide name from others. But original author can still see themselves?
        # The prompt says "whether anonymous or their name is visible in the community".
        # We'll show "Anonymous" for all if is_anon is true.
        show_author = "Anonymous" if is_anon else m.author.username
        show_initials = "??" if is_anon else m.author.username[:2].upper()
        show_color = "#9CA3AF" if is_anon else ('#8C1007' if m.author == request.user else '#4B5563')

        materials_list.append({
            'id':            m.id,
            'title':         m.title,
            'author':        show_author,
            'authorInitials': show_initials,
            'authorColor':   show_color,
            'is_anonymous':  is_anon,
            'is_mine':       m.author == request.user,
            'category':        m.category,
            'subject':       m.subject,
            'preview':       m.content,
            'likes':         m_likes,
            'views':         m.views,
            'comments':      m.comments.count(),
            'timeAgo':       'Just now',
            'emoji':         m.emoji,
            'liked':         m.likes.filter(id=request.user.id).exists(),
            'tags':          [m.subject],
            'file_url':      m.file.url if m.file else None,
        })

    # --- Top Contributors Logic ---
    all_users = User.objects.all()
    contributors = []
    
    for u in all_users:
        material_count = SharedMaterial.objects.filter(author=u).count()
        # Likes received by this user on their materials
        likes_received = SharedMaterial.objects.filter(author=u).aggregate(total=Count('likes'))['total'] or 0
        comment_count = Comment.objects.filter(author=u).count()
        completed_tasks = Task.objects.filter(user=u, completed=True).count()
        
        points = (material_count * 10) + (likes_received * 5) + (comment_count * 2) + (completed_tasks * 1)
        
        if points > 0:
            contributors.append({
                'username': u.username,
                'initials': u.username[:2].upper(),
                'points': points,
                'materials': material_count
            })
    
    # Sort by points descending
    contributors = sorted(contributors, key=lambda x: x['points'], reverse=True)[:5]

    return render(request, "main/collaborate.html", {
        'materials_json':        json.dumps(materials_list),
        'active_students':       User.objects.count(),
        'total_community_likes': total_community_likes,
        'top_contributors':      contributors,
    })


@login_required
@require_POST
@ratelimit(key='ip', rate='10/m', block=True)
def share_material(request):
    try:
        # Handle FormData (Multipart) correctly instead of reading the body stream
        title    = bleach.clean(request.POST.get('title', '').strip())
        subject  = bleach.clean(request.POST.get('subject', '').strip())
        category = bleach.clean(request.POST.get('category', 'General').strip())
        content  = bleach.clean(request.POST.get('preview', '').strip())
        is_anon  = request.POST.get('is_anonymous') == 'true'
        file_obj = request.FILES.get('file')

        if not title or len(title) > 255:
            return JsonResponse({'status': 'error', 'message': 'Title length invalid.'}, status=400)
        if len(subject) > 100:
            return JsonResponse({'status': 'error', 'message': 'Subject too long.'}, status=400)
        if len(category) > 20:
            return JsonResponse({'status': 'error', 'message': 'Category too long.'}, status=400)
        if not content:
            return JsonResponse({'status': 'error', 'message': 'Content missing.'}, status=400)

        material = SharedMaterial.objects.create(
            author   = request.user,
            title    = title,
            subject  = subject,
            category = category,
            content  = content,
            file     = file_obj,
            is_anonymous = is_anon,
            emoji    = '📄',
        )
        # Audit Log Fix
        log = AuditLog(user=request.user, action="Shared Material")
        log.details = f"Title: {title}"
        log.save()
        return JsonResponse({'status': 'success', 'material': {
            'id':            material.id,
            'title':         material.title,
            'author':        "Anonymous" if is_anon else material.author.username,
            'authorInitials': "??" if is_anon else material.author.username[:2].upper(),
            'authorColor':   "#9CA3AF" if is_anon else "#8C1007",
            'is_anonymous':  is_anon,
            'is_mine':       True,
            'category':        material.category,
            'subject':       material.subject,
            'preview':       material.content,
            'likes': 0, 'views': 0, 'comments': 0,
            'timeAgo': 'Just now',
            'emoji':  material.emoji,
            'liked':  False,
            'tags':   [material.subject],
            'file_url': material.file.url if material.file else None,
        }})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@login_required
@require_POST
def toggle_like_material(request, material_id):
    material = get_object_or_404(SharedMaterial, id=material_id)
    if material.likes.filter(id=request.user.id).exists():
        material.likes.remove(request.user)
        liked = False
    else:
        material.likes.add(request.user)
        liked = True
    return JsonResponse({'status': 'success', 'liked': liked, 'likes_count': material.likes.count()})


@login_required
def get_material_comments(request, material_id):
    material     = get_object_or_404(SharedMaterial, id=material_id)
    comments     = material.comments.all().order_by('-created_at')
    comments_list = [{
        'id':            c.id,
        'author':        c.author.username,
        'authorInitials': c.author.username[:2].upper(),
        'authorColor':   '#8C1007' if c.author == request.user else '#4B5563',
        'text':          c.text,
        'timeAgo':       'Just now',
    } for c in comments]
    return JsonResponse({'status': 'success', 'comments': comments_list})


@login_required
@require_POST
@ratelimit(key='ip', rate='20/m', block=True)
def add_comment(request, material_id):
    try:
        material = get_object_or_404(SharedMaterial, id=int(material_id))
        data     = json.loads(request.body)
        
        import bleach
        text     = bleach.clean(str(data.get('text', '')).strip())

        if not text:
            return JsonResponse({'status': 'error', 'message': 'Comment text required.'}, status=400)

        comment  = Comment.objects.create(material=material, author=request.user, text=text)
        log = AuditLog(user=request.user, action="Added Comment")
        log.details = f"Material ID: {material.id}"
        log.save()
        return JsonResponse({'status': 'success', 'comment': {
            'id':            comment.id,
            'author':        comment.author.username,
            'authorInitials': comment.author.username[:2].upper(),
            'authorColor':   '#8C1007',
            'text':          comment.text,
            'timeAgo':       'Just now',
        }})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)



@login_required
@csrf_protect
def progress(request):
    user_tasks      = Task.objects.filter(user=request.user)
    total_tasks     = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    summaries       = SummarizedDocument.objects.filter(user=request.user)
    today           = date.today()

    activity_dates = set()
    for t in user_tasks.filter(completed=True):
        activity_dates.add(t.created_at.date())
    for s in summaries:
        activity_dates.add(s.created_at.date())

    sorted_dates = sorted(activity_dates, reverse=True)
    streak       = 0
    if sorted_dates and sorted_dates[0] in (today, today - timedelta(days=1)):
        current_date = sorted_dates[0]
        streak       = 1
        for i in range(1, len(sorted_dates)):
            if sorted_dates[i] == current_date - timedelta(days=1):
                streak      += 1
                current_date = sorted_dates[i]
            else:
                break

    subjects = {}
    for t in user_tasks:
        sub = t.subject or 'General'
        subjects[sub] = subjects.get(sub, 0) + 1
    for s in summaries:
        sub = s.subject or 'General'
        subjects[sub] = subjects.get(sub, 0) + 1

    weekly_trend = []
    for i in range(3, -1, -1):
        start = today - timedelta(days=(i + 1) * 7)
        end   = today - timedelta(days=i * 7)
        weekly_trend.append(
            user_tasks.filter(completed=True, created_at__range=(start, end)).count() * 2
            + summaries.filter(created_at__range=(start, end)).count()
        )

    context = {
        'total_tasks':       total_tasks,
        'completed_count':   completed_tasks,
        'completion_rate':   round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0,
        'streak':            streak,
        'study_hours':       (completed_tasks * 2) + summaries.count(),
        'summaries_count':   summaries.count(),
        'category_stats_json': json.dumps([{'category': 'General Progress', 'completed': completed_tasks, 'total': total_tasks}]),
        'subject_labels_json': json.dumps(list(subjects.keys())),
        'subject_data_json':   json.dumps(list(subjects.values())),
        'weekly_hours_json':   json.dumps(weekly_trend),
    }
    return render(request, "main/progress.html", context)



@login_required
@csrf_protect
def profile(request):
    if request.method == "POST":
        new_username = request.POST.get('username', '').strip()
        
        # Validation
        if not new_username:
            messages.error(request, "Username cannot be empty.")
        elif new_username == request.user.username:
            messages.info(request, "No changes made.")
        elif User.objects.filter(username=new_username).exists():
            messages.error(request, "Username already taken.")
        elif not re.match(r'^[a-zA-Z0-9_\.\-]{3,150}$', new_username):
            messages.error(request, "Invalid username format. (3-150 chars, letters, numbers, dots, dashes, underscores)")
        else:
            old_un = request.user.username
            request.user.username = new_username
            request.user.save()
            log = AuditLog(user=request.user, action="Username Updated")
            log.details = f"Changed from {old_un} to {new_username}"
            log.save()
            messages.success(request, f"Username updated to {new_username} successfully!")
            return redirect('profile')

    profile_obj, _ = UserProfile.objects.get_or_create(user=request.user)

    user_tasks      = Task.objects.filter(user=request.user)
    total_tasks     = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    summaries_count = SummarizedDocument.objects.filter(user=request.user).count()

    context = {
        'user_level':          (completed_tasks // 5) + 1,
        'next_level_progress': int(((completed_tasks % 5) / 5) * 100) if (completed_tasks % 5) != 0 or completed_tasks == 0 else 100,
        'docs_count':          summaries_count,
        'completed_count':     completed_tasks,
        'total_tasks':         total_tasks,
        'completion_rate':     round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0,
        'study_hours':         (completed_tasks * 2) + summaries_count,
        'mfa_enabled':         profile_obj.totp_enabled,
        'recent_summaries':    SummarizedDocument.objects.filter(user=request.user).order_by('-created_at')[:5],
    }
    return render(request, "main/profile.html", context)


@login_required
@require_POST
def toggle_mfa(request):
    profile_obj, _ = UserProfile.objects.get_or_create(user=request.user)
    action = request.POST.get('action')
    
    if action == 'disable':
        profile_obj.totp_enabled = False
        profile_obj.save()
        send_security_alert(
            request.user, 
            "MFA Disabled", 
            "Multi-Factor Authentication (MFA) has been disabled for your account. Your account is now less secure against unauthorized access."
        )
        log = AuditLog(user=request.user, action="MFA Disabled")
        log.details = "User manually disabled TOTP MFA"
        log.save()
        messages.warning(request, "MFA has been disabled. We strongly recommend re-enabling it.")
    
    return redirect('profile')



@login_required
def download_summary_pdf(request, doc_id):
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    import io

    doc_obj = get_object_or_404(SummarizedDocument, id=doc_id, user=request.user)
    
    # We use io.BytesIO() as a buffer for the PDF
    buffer = io.BytesIO()
    
    # Create the SimpleDocTemplate with letter size
    final_pdf_doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    
    styles = getSampleStyleSheet()
    
    # Custom Title style
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#8C1007'),
        alignment=1, # Center
        spaceAfter=20,
        fontName='Helvetica-Bold'
    )
    
    # Body style
    body_style = ParagraphStyle(
        'BodyStyle',
        parent=styles['Normal'],
        fontSize=11,
        leading=14,
        fontName='Helvetica',
        alignment=4 # Justify
    )
    
    # Clean the title and summary for PDF compatibility (Remove only broken emojis, not accented characters)
    # This regex is specifically for emojis and non-standard symbols that break Helvetica
    emo_regex = r'[\U00010000-\U0010ffff]'
    clean_title = re.sub(emo_regex, '', os.path.splitext(doc_obj.file_name)[0]).strip()
    
    # Process text for bold and line breaks
    text_content = doc_obj.summary_text
    # Safely remove emojis
    text_content = re.sub(emo_regex, '', text_content)
    # Convert MarkDown style **bold** to ReportLab <b> HTML style
    text_content = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text_content)

    try:
        # Final safety cleanup for ReportLab (it only supports a limited subset of characters in Helvetica)
        def safe_text(txt):
            if not txt: return ""
            # Only keep characters that are likely to be in the standard WinAnsi encoding set
            # This is more inclusive than ASCII but avoids characters that crash Helvetica
            return re.sub(r'[^\x00-\xff\u2013\u2014\u2018\u2019\u201c\u201d\u2022]', '', txt)

        elements = []
        # Title with safety and bold emphasis
        elements.append(Paragraph(f"<b>Study Reviewer: {safe_text(clean_title)}</b>", title_style))
        elements.append(Spacer(1, 25))
        
        # Summary Content processing
        paragraphs = text_content.split('\n')
        for p_text in paragraphs:
            p_text = safe_text(p_text.strip())
            if p_text:
                # Basic cleaning for Paragraph (it hates raw < or > unless escaped)
                p_text = p_text.replace('<', '&lt;').replace('>', '&gt;')
                # Restore our bold tags which were escaped
                p_text = p_text.replace('&lt;b&gt;', '<b>').replace('&lt;/b&gt;', '</b>')
                
                # List detection (supporting common bullet points and numbers)
                if p_text.startswith('- ') or p_text.startswith('* ') or p_text.startswith('• ') or re.match(r'^\d+\.', p_text):
                    list_style = ParagraphStyle('ListStyle', parent=body_style, leftIndent=25, bulletIndent=10, spaceAfter=8)
                    elements.append(Paragraph(p_text, list_style))
                else:
                    elements.append(Paragraph(p_text, body_style))
                    elements.append(Spacer(1, 10))
        
        # Build the PDF
        final_pdf_doc.build(elements)
        
        pdf = buffer.getvalue()
        buffer.close()
        
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="Summary_{doc_id}.pdf"'
        response['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0'
        # Audit Secure Download
        log = AuditLog(user=request.user, action="Downloaded PDF Summary")
        log.details = f"Document ID: {doc_id}"
        log.save()
        return response

    except Exception as e:
        print(f"CRITICAL PDF GEN ERROR: {str(e)}")
        # If PDF fails, return as a simple text file so user still gets their content!
        response = HttpResponse(doc_obj.summary_text, content_type='text/plain')
        response['Content-Disposition'] = f'attachment; filename="Summary_{doc_id}.txt"'
        return response

@login_required
def download_shared_pdf(request, material_id):
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    import io

    material = get_object_or_404(SharedMaterial, id=material_id)
    
    buffer = io.BytesIO()
    final_pdf_doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=18)
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle('TitleStyle', parent=styles['Heading1'], fontSize=22, textColor=colors.HexColor('#8C1007'), alignment=1, spaceAfter=20, fontName='Helvetica-Bold')
    body_style = ParagraphStyle('BodyStyle', parent=styles['Normal'], fontSize=11, leading=14, fontName='Helvetica', alignment=4)
    
    emo_regex = r'[\U00010000-\U0010ffff]'
    clean_title = re.sub(emo_regex, '', material.title).strip()
    text_content = re.sub(emo_regex, '', material.content)
    text_content = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', text_content)

    try:
        def safe_text(txt):
            if not txt: return ""
            return re.sub(r'[^\x00-\xff\u2013\u2014\u2018\u2019\u201c\u201d\u2022]', '', txt)

        elements = []
        elements.append(Paragraph(f"<b>Community Resource: {safe_text(clean_title)}</b>", title_style))
        elements.append(Paragraph(f"Shared by: {material.author.username if not material.is_anonymous else 'Anonymous'}", styles['Italic']))
        elements.append(Spacer(1, 20))
        
        paragraphs = text_content.split('\n')
        for p_text in paragraphs:
            p_text = safe_text(p_text.strip())
            if p_text:
                p_text = p_text.replace('<', '&lt;').replace('>', '&gt;').replace('&lt;b&gt;', '<b>').replace('&lt;/b&gt;', '</b>')
                if p_text.startswith('- ') or p_text.startswith('* ') or re.match(r'^\d+\.', p_text):
                    list_style = ParagraphStyle('ListStyle', parent=body_style, leftIndent=25, bulletIndent=10, spaceAfter=8)
                    elements.append(Paragraph(p_text, list_style))
                else:
                    elements.append(Paragraph(p_text, body_style))
                    elements.append(Spacer(1, 10))
        
        final_pdf_doc.build(elements)
        pdf = buffer.getvalue()
        buffer.close()
        
        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="Shared_{material_id}.pdf"'
        return response
    except Exception as e:
        print(f"SHARED PDF GEN ERROR: {str(e)}")
        response = HttpResponse(material.content, content_type='text/plain')
        response['Content-Disposition'] = f'attachment; filename="Shared_{material_id}.txt"'
        return response

@login_required
@require_POST
def add_schedule_item(request):
    try:
        from datetime import datetime
        data = json.loads(request.body)
        date_str = data.get('date')
        date_obj = datetime.strptime(date_str, '%Y-%m-%d').date() if date_str else None
        
        item = ScheduleItem.objects.create(
            user     = request.user,
            day      = data.get('day', 'General'),
            date     = date_obj,
            time     = data.get('time'),
            activity = data.get('activity'),
            color    = data.get('color', 'blue'),
        )
        return JsonResponse({'status': 'success', 'item': {
            'id': item.id, 'day': item.day, 'date': item.date.strftime('%Y-%m-%d') if item.date else None,
            'time': item.time, 'activity': item.activity, 'color': item.color,
        }})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)


@login_required
@require_POST
@login_required
@require_POST
def delete_schedule_item(request, item_id):
    item = get_object_or_404(ScheduleItem, id=item_id, user=request.user)
    item.delete()
    return JsonResponse({'status': 'success'})

@login_required
@require_POST
def edit_schedule_item(request, item_id):
    try:
        from datetime import datetime
        item = get_object_or_404(ScheduleItem, id=item_id, user=request.user)
        data = json.loads(request.body)
        
        date_str = data.get('date')
        if date_str:
            item.date = datetime.strptime(date_str, '%Y-%m-%d').date()
            
        item.day      = data.get('day', item.day)
        item.time     = data.get('time', item.time)
        item.activity = data.get('activity', item.activity)
        item.color    = data.get('color', item.color)
        item.save()
        
        return JsonResponse({'status': 'success', 'item': {
            'id': item.id, 'day': item.day, 'date': item.date.strftime('%Y-%m-%d') if item.date else None,
            'time': item.time, 'activity': item.activity, 'color': item.color,
        }})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)